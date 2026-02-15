"""
FastAPI routes for Course Generation.
Clean, well-documented endpoints following REST conventions.
"""

from fastapi import APIRouter, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from typing import List, Optional
import asyncio
import json
import logging
import re
import time

from services.course_service import CourseService
from services.wetrocloud_youtube import WetroCloudYouTubeService
from clients.jina_reader_client import extract_web_content
from utils.file_storage import ProgressStorage
from clients.supabase_client import (
    upsert_chapter_progress,
    insert_course_quiz_attempt,
    get_chapter_progress_for_course,
)
from models.course_models import (
    LearningProfileRequest, ProgressUpdateRequest
)
from utils.exceptions import (
    NutonError, ValidationError, NotFoundError, GenerationError,
)

# File processing imports
from processors.mistral_ocr_extractor import MistralOCRExtractor, MistralOCRConfig
from prompts.course_prompts import build_topic_extraction_prompt
from utils.model_config import ModelConfig
import os
import tempfile

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/api/v1", tags=["courses"])

# Initialize services
course_service = CourseService()
yt_service = WetroCloudYouTubeService()


# Learning Profile Endpoints
@router.post("/learning-profile", status_code=201)
async def create_learning_profile(request: LearningProfileRequest):
    """
    Save or update user learning preferences.

    The 6 personalization questions:
    1. expertise: beginner/intermediate/advanced (default: beginner)
    2. format_pref: reading/listening/testing/mixed
    3. depth_pref: quick/detailed/conversational/academic
    4. role: student/professional/graduate_student
    5. learning_goal: exams/career/curiosity/supplement
    6. example_pref: real_world/technical/stories/analogies
    """
    success = course_service.save_learning_profile(request.dict())
    if not success:
        raise GenerationError("Failed to save profile", error_code="INTERNAL_ERROR")

    return {
        "success": True,
        "profile": request.dict(),
        "message": "Learning preferences saved"
    }


@router.get("/learning-profile/{user_id}")
async def get_learning_profile(user_id: str):
    """Retrieve user's learning preferences"""
    profile = course_service.get_learning_profile(user_id)
    
    if not profile:
        return {
            "exists": False,
            "message": "No profile found. User will use defaults."
        }
    
    return {
        "exists": True,
        "profile": profile
    }


# Course Generation Endpoints
@router.post("/courses/from-topic", status_code=201)
async def create_course_from_topic(
    user_id: str = Form(...),
    topic: str = Form(...),
    model: Optional[str] = Form(None),
    files: Optional[List[UploadFile]] = File(None),
    youtube_urls: Optional[str] = Form(None),
    web_urls: Optional[str] = Form(None)
):
    """
    Generate complete course from topic string with optional supplementary sources.

    **Blocking operation** - Returns full course in 30-60 seconds.

    Accepts Form data (not JSON) to support optional file uploads.
    Topic drives the course structure; files, YouTube videos, and web URLs
    provide supplementary context.

    Source params (all optional, can be combined):
    - files: PDF/PPTX/DOCX/TXT/MD/PNG/JPG/WEBP uploads (max 10, 50MB each)
    - youtube_urls: JSON array of YouTube URLs, e.g. '["https://youtube.com/watch?v=abc"]'
    - web_urls: JSON array of web URLs, e.g. '["https://example.com/article"]'

    Models available:
    - claude-haiku-4-5 (default): Fast, affordable, native web search
    - claude-sonnet-4-5: High quality, native web search
    - claude-opus-4-6: Highest quality, native web search
    - gpt-4o: Good quality, web search via Responses API
    - gpt-5-mini: Fast, affordable, web search via Responses API
    - gpt-5.2: High quality, web search via Responses API
    - llama-4-scout: Fastest, cheapest, Perplexity search fallback
    - llama-4-maverick: Fast, affordable, Perplexity search fallback
    """
    # Validate model if provided
    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(
            f"Invalid model. Available: {ModelConfig.get_available_models()}",
            error_code="INVALID_MODEL",
            context={"model": model},
        )

    # Parse URL lists
    try:
        yt_urls = _parse_url_list(youtube_urls)
        w_urls = _parse_url_list(web_urls)
    except ValueError as e:
        raise ValidationError(str(e), error_code="INVALID_URL_FORMAT")

    # Validate file count
    has_files = files and len(files) > 0 and files[0].filename
    if has_files and len(files) > 10:
        raise ValidationError("Maximum 10 files allowed", error_code="FILE_LIMIT_EXCEEDED")

    # Build merged source list (process all source types concurrently)
    tasks = []
    if has_files:
        tasks.append(_process_uploaded_files(files, model=model))
    if yt_urls:
        tasks.append(_process_youtube_urls(yt_urls, model=model))
    if w_urls:
        tasks.append(_process_web_urls(w_urls, model=model))

    all_processed = []
    if tasks:
        results = await asyncio.gather(*tasks)
        for result in results:
            all_processed.extend(result)

    # Route to appropriate service method
    if all_processed:
        result = await course_service.create_course_from_topic_with_files(
            user_id=user_id,
            topic=topic,
            files=all_processed,
            model=model
        )
    else:
        # Pure topic path
        result = await course_service.create_course_from_topic(
            user_id=user_id,
            topic=topic,
            context={},
            model=model
        )

    return result


@router.post("/courses/from-files", status_code=201)
async def create_course_from_files(
    files: Optional[List[UploadFile]] = File(None),
    user_id: str = Form(...),
    organization: str = Form("auto"),
    model: Optional[str] = Form(None),
    youtube_urls: Optional[str] = Form(None),
    web_urls: Optional[str] = Form(None)
):
    """
    Generate course from uploaded files, YouTube videos, and/or web URLs.

    At least one source is required (file, YouTube URL, or web URL).

    **Multi-source support:**
    - Files: PDF/PPTX/DOCX/TXT/MD/PNG/JPG/WEBP (max 10, 50MB each)
    - youtube_urls: JSON array of YouTube URLs, e.g. '["https://youtube.com/watch?v=abc"]'
    - web_urls: JSON array of web URLs, e.g. '["https://example.com/article"]'

    Organization options:
    - auto: Let system decide based on topic similarity
    - thematic_bridge: Unified course showing connections
    - sequential_sections: Separate sections within one course
    - separate_courses: Create multiple independent courses

    Response format varies by organization:
    - Single course: `{course_id, status, course, generation_time_seconds, ...}`
    - Separate courses: `{organization: "separate_courses", total_courses, courses: [...]}`
    """
    # Validate model
    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(
            f"Invalid model. Available: {ModelConfig.get_available_models()}",
            error_code="INVALID_MODEL",
            context={"model": model},
        )

    # Parse URL lists
    try:
        yt_urls = _parse_url_list(youtube_urls)
        w_urls = _parse_url_list(web_urls)
    except ValueError as e:
        raise ValidationError(str(e), error_code="INVALID_URL_FORMAT")

    # Validate file count
    has_files = files and len(files) > 0 and files[0].filename
    if has_files and len(files) > 10:
        raise ValidationError("Maximum 10 files allowed", error_code="FILE_LIMIT_EXCEEDED")

    # Require at least one source
    if not has_files and not yt_urls and not w_urls:
        raise ValidationError(
            "At least 1 source required (file, YouTube URL, or web URL)",
            error_code="MISSING_SOURCE",
        )

    # Build merged source list (process all source types concurrently)
    tasks = []
    if has_files:
        tasks.append(_process_uploaded_files(files, model=model))
    if yt_urls:
        tasks.append(_process_youtube_urls(yt_urls, model=model))
    if w_urls:
        tasks.append(_process_web_urls(w_urls, model=model))

    all_processed = []
    if tasks:
        results = await asyncio.gather(*tasks)
        for result in results:
            all_processed.extend(result)

    if not all_processed:
        raise ValidationError(
            "No usable content extracted from provided sources",
            error_code="NO_CONTENT_EXTRACTED",
        )

    # Generate course
    result = await course_service.create_course_from_files(
        user_id=user_id,
        files=all_processed,
        organization=organization,
        model=model
    )

    return result


# SSE Streaming Helpers

def _sse_event(data: dict) -> str:
    """Format a dict as an SSE data line."""
    return f"data: {json.dumps(data, default=str)}\n\n"


def _sse_error_event(
    exc: Exception,
    phase: str = "stream",
    error_code: str = "INTERNAL_ERROR",
    status_code: int = 500,
    context: Optional[dict] = None,
) -> str:
    """Build a structured SSE error event from an exception."""
    if isinstance(exc, NutonError):
        error_code = exc.error_code
        status_code = exc.status_code
        context = exc.context
        message = exc.message
    else:
        message = str(exc)
    return _sse_event({
        "type": "error",
        "error": error_code,
        "message": message,
        "status_code": status_code,
        "phase": phase,
        "context": context,
    })


# SSE Streaming Endpoints

@router.post("/courses/from-topic/stream")
async def create_course_from_topic_stream(
    user_id: str = Form(...),
    topic: str = Form(...),
    model: Optional[str] = Form(None),
    files: Optional[List[UploadFile]] = File(None),
    youtube_urls: Optional[str] = Form(None),
    web_urls: Optional[str] = Form(None)
):
    """
    SSE streaming version of /courses/from-topic.
    Returns a text/event-stream that progressively emits:
    - processing_sources: when source extraction begins
    - outline_ready: course outline + metadata (~10s)
    - chapter_ready: each chapter as it finishes
    - course_complete: all done with timing
    - error: on failure
    """
    # Pre-stream validation (raises before stream starts → global handler returns JSON)
    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(
            f"Invalid model. Available: {ModelConfig.get_available_models()}",
            error_code="INVALID_MODEL",
            context={"model": model},
        )

    try:
        yt_urls = _parse_url_list(youtube_urls)
        w_urls = _parse_url_list(web_urls)
    except ValueError as e:
        raise ValidationError(str(e), error_code="INVALID_URL_FORMAT")

    has_files = files and len(files) > 0 and files[0].filename
    if has_files and len(files) > 10:
        raise ValidationError("Maximum 10 files allowed", error_code="FILE_LIMIT_EXCEEDED")

    async def event_generator():
        try:
            # Process sources
            all_processed = []
            source_count = 0

            if has_files:
                source_count += len(files)
            if yt_urls:
                source_count += len(yt_urls)
            if w_urls:
                source_count += len(w_urls)

            if source_count > 0:
                yield _sse_event({
                    "type": "processing_sources",
                    "message": f"Processing {source_count} source(s)...",
                    "source_count": source_count
                })

            # Process all source types concurrently
            tasks = []
            if has_files:
                tasks.append(_process_uploaded_files(files, model=model))
            if yt_urls:
                tasks.append(_process_youtube_urls(yt_urls, model=model))
            if w_urls:
                tasks.append(_process_web_urls(w_urls, model=model))

            if tasks:
                results = await asyncio.gather(*tasks)
                for result in results:
                    all_processed.extend(result)

            # Route to streaming service method
            if all_processed:
                stream = course_service.create_course_from_topic_with_files_stream(
                    user_id=user_id,
                    topic=topic,
                    files=all_processed,
                    model=model
                )
            else:
                stream = course_service.create_course_from_topic_stream(
                    user_id=user_id,
                    topic=topic,
                    context={},
                    model=model
                )

            async for event in stream:
                yield _sse_event(event)

        except Exception as e:
            logger.error(f"SSE stream error: {e}")
            yield _sse_error_event(e, phase="stream")

    return StreamingResponse(event_generator(), media_type="text/event-stream")


@router.post("/courses/from-files/stream")
async def create_course_from_files_stream(
    files: Optional[List[UploadFile]] = File(None),
    user_id: str = Form(...),
    organization: str = Form("auto"),
    model: Optional[str] = Form(None),
    youtube_urls: Optional[str] = Form(None),
    web_urls: Optional[str] = Form(None)
):
    """
    SSE streaming version of /courses/from-files.
    Returns a text/event-stream that progressively emits events.
    """
    # Pre-stream validation
    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(
            f"Invalid model. Available: {ModelConfig.get_available_models()}",
            error_code="INVALID_MODEL",
            context={"model": model},
        )

    try:
        yt_urls = _parse_url_list(youtube_urls)
        w_urls = _parse_url_list(web_urls)
    except ValueError as e:
        raise ValidationError(str(e), error_code="INVALID_URL_FORMAT")

    has_files = files and len(files) > 0 and files[0].filename
    if has_files and len(files) > 10:
        raise ValidationError("Maximum 10 files allowed", error_code="FILE_LIMIT_EXCEEDED")

    if not has_files and not yt_urls and not w_urls:
        raise ValidationError(
            "At least 1 source required (file, YouTube URL, or web URL)",
            error_code="MISSING_SOURCE",
        )

    async def event_generator():
        try:
            # Process sources
            all_processed = []
            source_count = (len(files) if has_files else 0) + len(yt_urls) + len(w_urls)

            yield _sse_event({
                "type": "processing_sources",
                "message": f"Processing {source_count} source(s)...",
                "source_count": source_count
            })

            # Process all source types concurrently
            tasks = []
            if has_files:
                tasks.append(_process_uploaded_files(files, model=model))
            if yt_urls:
                tasks.append(_process_youtube_urls(yt_urls, model=model))
            if w_urls:
                tasks.append(_process_web_urls(w_urls, model=model))

            if tasks:
                results = await asyncio.gather(*tasks)
                for result in results:
                    all_processed.extend(result)

            if not all_processed:
                yield _sse_error_event(
                    ValidationError("No usable content extracted from provided sources", error_code="NO_CONTENT_EXTRACTED"),
                    phase="source_processing",
                )
                return

            # Stream course generation
            async for event in course_service.create_course_from_files_stream(
                user_id=user_id,
                files=all_processed,
                organization=organization,
                model=model
            ):
                yield _sse_event(event)

        except Exception as e:
            logger.error(f"SSE stream error: {e}")
            yield _sse_error_event(e, phase="stream")

    return StreamingResponse(event_generator(), media_type="text/event-stream")


# Course Access Endpoints
@router.get("/courses/{course_id}")
async def get_course(course_id: str):
    """
    Retrieve full course with all chapters.
    
    Returns complete course data including:
    - Course metadata
    - Chapter summaries
    - Progress information (if available)
    """
    course = course_service.get_course(course_id)

    if not course:
        raise NotFoundError("Course not found", error_code="COURSE_NOT_FOUND", context={"course_id": course_id})

    # Get progress if available
    progress_storage = ProgressStorage()
    progress = progress_storage.load_progress(course["user_id"], course_id)

    return {
        "course": course,
        "progress": progress
    }


@router.post("/courses/{course_id}/resume")
async def resume_course(course_id: str, model: Optional[str] = None):
    """
    Resume generation of a partially-completed course.
    Detects missing or errored chapters and generates only the gaps.
    Also generates study guide and flashcards if missing.
    Re-entrant: safe to call multiple times.

    Query params:
        model: Optional model override (defaults to course's original model)
    """
    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(
            f"Invalid model. Available: {ModelConfig.get_available_models()}",
            error_code="INVALID_MODEL",
            context={"model": model},
        )

    result = await course_service.resume_course(
        course_id=course_id,
        model=model,
    )
    return result


@router.post("/courses/{course_id}/resume/stream")
async def resume_course_stream(course_id: str, model: Optional[str] = None):
    """
    SSE streaming version of course resume.
    Emits: resume_started, chapter_ready, study_guide_ready, flashcards_ready, course_complete, error
    """
    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(
            f"Invalid model. Available: {ModelConfig.get_available_models()}",
            error_code="INVALID_MODEL",
            context={"model": model},
        )

    async def event_generator():
        try:
            async for event in course_service.resume_course_stream(
                course_id=course_id,
                model=model,
            ):
                yield _sse_event(event)
                if event.get("type") == "error":
                    return
        except Exception as e:
            logger.error(f"Resume SSE stream error: {e}")
            yield _sse_error_event(e, phase="resume")

    return StreamingResponse(event_generator(), media_type="text/event-stream")


@router.get("/courses/by-slug/{slug}")
async def get_course_by_slug(slug: str):
    """
    Retrieve full course by its URL-friendly slug.
    e.g. GET /api/v1/courses/by-slug/exploring-modern-ai
    """
    from clients.supabase_client import get_course_by_slug as db_get_course_by_slug, get_chapters_by_course

    course = db_get_course_by_slug(slug)
    if not course:
        raise NotFoundError("Course not found", error_code="COURSE_NOT_FOUND", context={"slug": slug})

    # Attach chapters
    chapters = get_chapters_by_course(course["id"])
    course["chapters"] = chapters

    # Get progress if available
    progress_storage = ProgressStorage()
    progress = progress_storage.load_progress(course["user_id"], course["id"])

    return {
        "course": course,
        "progress": progress
    }


@router.get("/courses/{course_id}/chapters/{chapter_order}")
async def get_chapter(course_id: str, chapter_order: int):
    """
    Retrieve specific chapter by order number (1-indexed).

    Returns full chapter content including:
    - Markdown content
    - Quiz questions
    - Source citations
    """
    chapter = course_service.get_chapter(course_id, chapter_order)

    if not chapter:
        raise NotFoundError("Chapter not found", error_code="CHAPTER_NOT_FOUND", context={"course_id": course_id, "chapter_order": chapter_order})

    return {"chapter": chapter}


@router.get("/users/{user_id}/courses")
async def list_user_courses(user_id: str):
    """List all courses for a user"""
    courses = course_service.list_user_courses(user_id)
    return {
        "user_id": user_id,
        "course_count": len(courses),
        "courses": courses
    }


# Progress Endpoints
@router.post("/courses/{course_id}/progress")
async def update_progress(course_id: str, request: ProgressUpdateRequest):
    """
    Update chapter completion and quiz scores.
    
    Track:
    - Chapter completion status
    - Quiz scores
    - Time spent
    """
    course = course_service.get_course(course_id)
    if not course:
        raise NotFoundError("Course not found", error_code="COURSE_NOT_FOUND", context={"course_id": course_id})

    user_id = course["user_id"]
    now = time.strftime("%Y-%m-%dT%H:%M:%SZ")

    # Upsert chapter progress directly to Supabase
    upsert_chapter_progress(
        user_id=user_id,
        course_id=course_id,
        chapter_id=request.chapter_id,
        completed=request.completed,
        time_spent_minutes=request.time_spent_minutes or 0,
        completed_at=now if request.completed else None,
    )

    # Insert quiz attempt if score provided
    if request.quiz_score is not None:
        insert_course_quiz_attempt(
            user_id=user_id,
            chapter_id=request.chapter_id,
            score=request.quiz_score,
            completed_at=now,
        )

    # Re-query to compute overall stats
    rows = get_chapter_progress_for_course(user_id, course_id)
    total = course["total_chapters"]
    completed_count = sum(1 for r in rows if r.get("completed"))
    percentage = round((completed_count / total) * 100) if total > 0 else 0

    return {
        "success": True,
        "overall_progress": {
            "completed_chapters": completed_count,
            "total_chapters": total,
            "percentage": percentage,
            "last_activity": now,
        }
    }


@router.get("/users/{user_id}/progress")
async def get_user_progress(user_id: str):
    """Get learning progress across all courses (batched — 3 queries total)"""
    courses = course_service.list_user_courses(user_id)

    # Batch: 2 queries for ALL progress + quiz attempts (instead of N+1+N*M)
    progress_storage = ProgressStorage()
    progress_map = progress_storage.load_all_progress(user_id, courses)

    courses_with_progress = []
    for course in courses:
        courses_with_progress.append({
            "course": course,
            "progress": progress_map.get(course["id"])
        })

    # Calculate stats
    completed_courses = [c for c in courses_with_progress if c["progress"] and c["progress"]["overall_progress"]["percentage"] == 100]
    in_progress = [c for c in courses_with_progress if c["progress"] and c["progress"]["overall_progress"]["percentage"] > 0 and c["progress"]["overall_progress"]["percentage"] < 100]

    return {
        "user_id": user_id,
        "total_courses": len(courses),
        "completed_courses": len(completed_courses),
        "in_progress": len(in_progress),
        "courses": courses_with_progress
    }


# Course Q&A Endpoint (with chat history)
@router.post("/courses/{course_id}/ask")
async def ask_course_question(
    course_id: str,
    question: str = Form(...),
    model: Optional[str] = Form(None),
    user_id: Optional[str] = Form(None)
):
    """
    Ask a question about a course. Answers are grounded in the uploaded source material.
    Uses Pinecone vector search to find relevant chunks from the original documents.
    When user_id is provided, maintains persistent chat history (Redis cache + Supabase).
    """
    result = await course_service.query_course(
        course_id=course_id,
        question=question,
        model=model,
        user_id=user_id
    )
    return result


# Notes Generation Endpoints
@router.post("/courses/{course_id}/generate-notes", status_code=201)
async def generate_notes(
    course_id: str,
    user_id: str = Form(...),
    model: Optional[str] = Form(None)
):
    """
    Generate comprehensive study notes from course source materials.
    Stores notes in courses.summary_md. Also updates course title/topic/slug
    from document analysis and generates a study guide.

    This is a blocking operation that may take 30-120s depending on material size.
    """
    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(
            f"Invalid model. Available: {ModelConfig.get_available_models()}",
            error_code="INVALID_MODEL",
            context={"model": model},
        )

    result = await course_service.generate_notes(
        course_id=course_id,
        user_id=user_id,
        model=model
    )
    return result


@router.post("/courses/{course_id}/generate-notes-flashcards", status_code=201)
async def generate_notes_flashcards(
    course_id: str,
    user_id: str = Form(...),
    model: Optional[str] = Form(None)
):
    """
    Generate flashcards from course source materials.
    Standalone operation — does not require generate-notes to run first.
    Saves to both courses.flashcards JSONB and flashcard_sets table.

    This is a blocking operation that may take 30-90s depending on material size.
    """
    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(
            f"Invalid model. Available: {ModelConfig.get_available_models()}",
            error_code="INVALID_MODEL",
            context={"model": model},
        )

    result = await course_service.generate_notes_flashcards(
        course_id=course_id,
        user_id=user_id,
        model=model
    )
    return result


@router.post("/courses/{course_id}/generate-notes-quiz", status_code=201)
async def generate_notes_quiz(
    course_id: str,
    user_id: str = Form(...),
    model: Optional[str] = Form(None)
):
    """
    Generate a comprehensive quiz from course source materials.
    Standalone operation — does not require generate-notes to run first.
    Saves to quiz_sets table with MCQ, fill-in-gap, and scenario questions.

    This is a blocking operation that may take 30-90s depending on material size.
    """
    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(
            f"Invalid model. Available: {ModelConfig.get_available_models()}",
            error_code="INVALID_MODEL",
            context={"model": model},
        )

    result = await course_service.generate_notes_quiz(
        course_id=course_id,
        user_id=user_id,
        model=model
    )
    return result


@router.get("/courses/{course_id}/notes-flashcards")
async def get_notes_flashcards(course_id: str):
    """
    Retrieve incrementally-inserted flashcards from flashcard_sets table.
    Returns partial results while generation is still running, or full set after.
    """
    from clients.supabase_client import get_course_flashcard_sets

    sets = get_course_flashcard_sets(course_id)
    if not sets:
        raise NotFoundError("No flashcard sets found for this course. Generate them first via POST /generate-notes-flashcards.", error_code="FLASHCARDS_NOT_FOUND", context={"course_id": course_id})

    all_cards = []
    for s in sets:
        cards = s.get("flashcards", [])
        if isinstance(cards, list):
            all_cards.extend(cards)

    return {
        "course_id": course_id,
        "total": len(all_cards),
        "sets_count": len(sets),
        "flashcards": all_cards
    }


@router.get("/courses/{course_id}/notes-quiz")
async def get_notes_quiz(course_id: str):
    """
    Retrieve incrementally-inserted quiz questions from quiz_sets table.
    Returns partial results while generation is still running, or full set after.
    """
    from clients.supabase_client import get_course_quiz_sets

    sets = get_course_quiz_sets(course_id)
    if not sets:
        raise NotFoundError("No quiz sets found for this course. Generate them first via POST /generate-notes-quiz.", error_code="QUIZ_NOT_FOUND", context={"course_id": course_id})

    merged = {"mcq": [], "fill_in_gap": [], "scenario": []}
    for s in sets:
        quiz = s.get("quiz", {})
        if isinstance(quiz, dict):
            for q_type in ("mcq", "fill_in_gap", "scenario"):
                merged[q_type].extend(quiz.get(q_type, []))

    total = sum(len(v) for v in merged.values())
    by_type = {k: len(v) for k, v in merged.items()}

    return {
        "course_id": course_id,
        "total_questions": total,
        "sets_count": len(sets),
        "by_type": by_type,
        "quiz": merged
    }


@router.get("/courses/{course_id}/notes")
async def get_course_notes(course_id: str):
    """
    Retrieve stored notes (summary_md) for a course.
    Returns 404 if notes have not been generated yet.
    """
    from clients.supabase_client import get_course_by_id

    course = get_course_by_id(course_id)
    if not course:
        raise NotFoundError("Course not found", error_code="COURSE_NOT_FOUND", context={"course_id": course_id})

    summary_md = course.get("summary_md")
    if not summary_md:
        raise NotFoundError("Notes not yet generated. Call POST /generate-notes first.", error_code="NOTES_NOT_FOUND", context={"course_id": course_id})

    return {
        "course_id": course_id,
        "title": course.get("title", ""),
        "notes_length": len(summary_md),
        "summary_md": summary_md
    }


# Study Guide Endpoint
@router.get("/courses/{course_id}/study-guide")
async def get_study_guide(course_id: str):
    """Get the study guide for a course."""
    from utils.file_storage import StudyGuideStorage
    study_guide = StudyGuideStorage.get_study_guide(course_id)
    if not study_guide:
        raise NotFoundError("Study guide not found. It may still be generating.", error_code="STUDY_GUIDE_NOT_FOUND", context={"course_id": course_id})
    return {"course_id": course_id, "study_guide": study_guide}


# Flashcards Endpoint
@router.get("/courses/{course_id}/flashcards")
async def get_flashcards(course_id: str):
    """Get the flashcards for a course."""
    from utils.file_storage import FlashcardStorage
    flashcards = FlashcardStorage.get_flashcards(course_id)
    if not flashcards:
        raise NotFoundError("Flashcards not found. They may still be generating.", error_code="FLASHCARDS_NOT_FOUND", context={"course_id": course_id})
    return {"course_id": course_id, "total": len(flashcards), "flashcards": flashcards}


@router.get("/courses/{course_id}/chapters/{chapter_order}/flashcards")
async def get_chapter_flashcards(course_id: str, chapter_order: int):
    """Get flashcards for a specific chapter."""
    from clients.supabase_client import get_chapter_by_order
    chapter = get_chapter_by_order(course_id, chapter_order)
    if not chapter:
        raise NotFoundError("Chapter not found", error_code="CHAPTER_NOT_FOUND", context={"course_id": course_id, "chapter_order": chapter_order})
    flashcards = chapter.get("flashcards") or []
    return {
        "course_id": course_id,
        "chapter_order": chapter_order,
        "chapter_title": chapter.get("title", ""),
        "total": len(flashcards),
        "flashcards": flashcards
    }


# Final Exam Endpoints
@router.post("/courses/{course_id}/generate-exam")
async def generate_exam(
    course_id: str,
    user_id: str = Form(...),
    exam_size: int = Form(30),
    model: Optional[str] = Form(None)
):
    """
    Generate a final exam for a course (on-demand).
    exam_size: 30 (15 MCQ / 8 fill-in / 7 theory) or 50 (25 MCQ / 15 fill-in / 10 theory)
    """
    if exam_size not in (30, 50):
        raise ValidationError("exam_size must be 30 or 50", error_code="INVALID_EXAM_SIZE", context={"exam_size": exam_size})

    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(f"Invalid model. Available: {ModelConfig.get_available_models()}", error_code="INVALID_MODEL", context={"model": model})

    result = await course_service.generate_final_exam(
        course_id=course_id,
        user_id=user_id,
        exam_size=exam_size,
        model=model
    )
    return result


@router.get("/courses/{course_id}/exam")
async def get_exam(course_id: str, user_id: Optional[str] = None):
    """Get the most recent exam for a course (optionally filtered by user)."""
    from utils.file_storage import ExamStorage
    exam = ExamStorage.get_exam(course_id, user_id)
    if not exam:
        raise NotFoundError("No exam found. Generate one first.", error_code="EXAM_NOT_FOUND", context={"course_id": course_id})
    return exam


# Exam Submission & Grading Endpoints
@router.post("/courses/{course_id}/exam/{exam_id}/submit")
async def submit_exam(
    course_id: str,
    exam_id: str,
    user_id: str = Form(...),
    answers: str = Form(...),
    model: Optional[str] = Form(None),
    time_taken_seconds: Optional[int] = Form(None)
):
    """
    Submit exam answers for grading.

    answers: JSON string with shape {"mcq": [2,0,3,...], "fill_in_gap": ["answer1",...], "theory": ["essay text",...]}
    - mcq: array of selected option indices (0-3)
    - fill_in_gap: array of answer strings
    - theory: array of essay answer strings

    Returns graded results with per-question breakdown and rubric analysis for theory.
    """
    # Parse answers JSON
    try:
        parsed_answers = json.loads(answers)
    except json.JSONDecodeError:
        raise ValidationError("answers must be valid JSON", error_code="INVALID_JSON")

    # Validate answer structure
    for key in ("mcq", "fill_in_gap", "theory"):
        if key not in parsed_answers:
            raise ValidationError(f"answers must contain '{key}' array", error_code="INVALID_ANSWERS_FORMAT")
        if not isinstance(parsed_answers[key], list):
            raise ValidationError(f"answers.{key} must be an array", error_code="INVALID_ANSWERS_FORMAT")

    if model and model not in ModelConfig.get_available_models():
        raise ValidationError(f"Invalid model. Available: {ModelConfig.get_available_models()}", error_code="INVALID_MODEL", context={"model": model})

    result = await course_service.grade_exam_submission(
        exam_id=exam_id,
        user_id=user_id,
        answers=parsed_answers,
        model=model,
        time_taken_seconds=time_taken_seconds
    )
    return result


@router.get("/courses/{course_id}/exam/{exam_id}/attempts")
async def get_exam_attempts(course_id: str, exam_id: str, user_id: Optional[str] = None):
    """
    Get all attempts for an exam by a specific user.
    Returns a list of attempts with scores and timestamps (without full results for brevity).
    """
    if not user_id:
        raise ValidationError("user_id query parameter is required", error_code="MISSING_REQUIRED_FIELD")

    from utils.file_storage import ExamAttemptStorage
    attempts = ExamAttemptStorage.get_attempts(exam_id, user_id)
    return {
        "exam_id": exam_id,
        "user_id": user_id,
        "total_attempts": len(attempts),
        "attempts": attempts
    }


@router.get("/courses/{course_id}/exam/{exam_id}/attempts/{attempt_id}")
async def get_exam_attempt_detail(course_id: str, exam_id: str, attempt_id: str):
    """
    Get full details of a specific exam attempt including rubric breakdowns.
    """
    from utils.file_storage import ExamAttemptStorage
    attempt = ExamAttemptStorage.get_attempt_by_id(attempt_id)
    if not attempt:
        raise NotFoundError("Attempt not found", error_code="ATTEMPT_NOT_FOUND", context={"attempt_id": attempt_id})
    return attempt


# Chat History Endpoints
@router.get("/courses/{course_id}/chat-history")
async def get_chat_history(course_id: str, user_id: Optional[str] = None, limit: int = 20):
    """Get chat history for a course+user pair."""
    if not user_id:
        raise ValidationError("user_id is required", error_code="MISSING_REQUIRED_FIELD")
    from utils.file_storage import ChatStorage
    messages = ChatStorage.get_messages(course_id, user_id, limit=limit)
    return {"course_id": course_id, "user_id": user_id, "messages": messages}


@router.delete("/courses/{course_id}/chat-history")
async def clear_chat_history(course_id: str, user_id: Optional[str] = None):
    """Clear chat history for a course+user pair."""
    if not user_id:
        raise ValidationError("user_id is required", error_code="MISSING_REQUIRED_FIELD")
    from utils.file_storage import ChatStorage
    ChatStorage.clear_messages(course_id, user_id)
    # Also clear Redis cache
    try:
        from clients.redis_client import clear_chat
        await clear_chat(course_id, user_id)
    except Exception:
        pass
    return {"success": True, "message": "Chat history cleared"}


# Helper functions
ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md", ".png", ".jpg", ".jpeg", ".webp"}
MAX_FILE_SIZE_MB = 50


async def _process_uploaded_files(files: List[UploadFile], model: Optional[str] = None) -> List[dict]:
    """Process uploaded files with OCR and topic extraction (parallel).
    Also fires off S3 uploads in the background so source files are persisted.

    Extraction strategy by format:
    - .txt/.md          → plaintext read (UTF-8)
    - .pdf              → Mistral OCR (with legacy fallback)
    - .png/.jpg/.jpeg/.webp → Mistral OCR image processing
    - .docx/.pptx       → S3 upload (blocking) → Jina Reader → python-docx/pptx fallback
    """
    from clients.s3_client import (
        build_s3_key, get_s3_url, get_s3_presigned_url,
        fire_and_forget_upload, upload_bytes_to_s3_async, get_content_type,
    )
    from utils.file_storage import generate_uuid as gen_uuid

    # Initialize Mistral OCR (for PDFs and images)
    mistral_config = MistralOCRConfig(
        enhance_metadata_with_llm=True,
        fallback_method="legacy",
        include_images=True,
    )
    extractor = MistralOCRExtractor(config=mistral_config)
    semaphore = asyncio.Semaphore(3)

    PLAINTEXT_EXTS = {".txt", ".md"}
    OFFICE_EXTS = {".docx", ".pptx"}
    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}

    # Validate all files first (instant, sequential)
    file_data = []
    for file in files:
        ext = os.path.splitext(file.filename)[1].lower() if file.filename else ""
        if ext not in ALLOWED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {ext}. Allowed: {ALLOWED_EXTENSIONS}")

        content = await file.read()
        await file.seek(0)

        size_mb = len(content) / (1024 * 1024)
        if size_mb > MAX_FILE_SIZE_MB:
            raise ValueError(f"File {file.filename} is {size_mb:.1f}MB. Max: {MAX_FILE_SIZE_MB}MB")

        file_data.append((file, ext, content))

    async def _process_single_file(file, ext, content):
        async with semaphore:
            file_id = gen_uuid()
            s3_key = build_s3_key(file_id, file.filename)
            s3_url = get_s3_url(s3_key)
            content_type = get_content_type(file.filename)

            # --- Plaintext ---
            if ext in PLAINTEXT_EXTS:
                fire_and_forget_upload(s3_key, content, content_type)
                text = content.decode("utf-8", errors="replace")
                if not text.strip():
                    raise ValueError(f"No text extracted from {file.filename}")
                topic = await _extract_topic(text[:2000], model=model)
                logger.info(f"Processed {file.filename}: {topic} ({len(text)} chars, plaintext)")
                return {
                    "filename": file.filename, "topic": topic,
                    "extracted_text": text, "pages": 1, "char_count": len(text),
                    "source_url": s3_url, "source_type": "text",
                }

            # --- Office documents (DOCX/PPTX) via Jina Reader ---
            if ext in OFFICE_EXTS:
                # Blocking S3 upload — we need the URL before calling Jina
                uploaded = await upload_bytes_to_s3_async(s3_key, content, content_type)
                if not uploaded:
                    raise ValueError(f"S3 upload failed for {file.filename}")

                # Presigned URL for Jina (temporary access, 15 min)
                presigned_url = get_s3_presigned_url(s3_key)
                text = await _extract_office_via_jina(presigned_url, file.filename, ext, content)
                topic = await _extract_topic(text[:2000], model=model)
                source_type = "pptx" if ext == ".pptx" else "docx"
                logger.info(f"Processed {file.filename}: {topic} ({len(text)} chars, {source_type})")
                return {
                    "filename": file.filename, "topic": topic,
                    "extracted_text": text, "pages": 0, "char_count": len(text),
                    "source_url": s3_url, "source_type": source_type,
                }

            # --- Images via Mistral OCR ---
            if ext in IMAGE_EXTS:
                fire_and_forget_upload(s3_key, content, content_type)
                temp_path = _save_temp_file(file, content)
                try:
                    extraction = await asyncio.to_thread(extractor.process_document, temp_path)
                    text = extraction.get('full_text', '')
                    if not text:
                        raise ValueError(f"No text extracted from {file.filename}")
                    topic = await _extract_topic(text[:2000], model=model)
                    logger.info(f"Processed {file.filename}: {topic} ({len(text)} chars, image OCR)")
                    return {
                        "filename": file.filename, "topic": topic,
                        "extracted_text": text, "pages": 1, "char_count": len(text),
                        "source_url": s3_url, "source_type": "image",
                    }
                finally:
                    if os.path.exists(temp_path):
                        os.remove(temp_path)

            # --- PDF via Mistral OCR (existing path) ---
            fire_and_forget_upload(s3_key, content, content_type)
            temp_path = _save_temp_file(file, content)
            try:
                extraction = await asyncio.to_thread(extractor.process_document, temp_path)
                text = extraction.get('full_text', '')
                if not text:
                    raise ValueError(f"No text extracted from {file.filename}")
                topic = await _extract_topic(text[:2000], model=model)
                logger.info(f"Processed {file.filename}: {topic} ({len(text)} chars, {extraction.get('total_pages', 0)} pages)")
                return {
                    "filename": file.filename, "topic": topic,
                    "extracted_text": text,
                    "pages": extraction.get('total_pages', 0), "char_count": len(text),
                    "source_url": s3_url, "source_type": "pdf",
                }
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)

    results = await asyncio.gather(*[
        _process_single_file(f, ext, content) for f, ext, content in file_data
    ])
    return list(results)


async def _extract_office_via_jina(s3_url: str, filename: str, ext: str, content: bytes) -> str:
    """Extract text from office documents via Jina Reader (S3 URL), with native fallback.

    Primary:  Jina Reader API via S3 URL
    Fallback: python-docx (for .docx) or python-pptx (for .pptx)
    """
    # Primary: Jina Reader
    try:
        logger.info(f"Extracting {filename} via Jina Reader: {s3_url}")
        result = await asyncio.to_thread(extract_web_content, s3_url, 60)
        if result.get("success") and result.get("text", "").strip():
            text = result["text"]
            logger.info(f"Jina Reader extracted {len(text)} chars from {filename}")
            return text
        logger.warning(f"Jina Reader returned no content for {filename}, trying native fallback")
    except Exception as e:
        logger.warning(f"Jina Reader failed for {filename}: {e}, trying native fallback")

    # Fallback: python-docx / python-pptx
    return _extract_office_native(filename, ext, content)


def _extract_office_native(filename: str, ext: str, content: bytes) -> str:
    """Fallback extraction using python-docx or python-pptx."""
    import io

    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)
            if not text.strip():
                raise ValueError(f"No text extracted from {filename}")
            logger.info(f"python-docx extracted {len(text)} chars from {filename}")
            return text
        except ImportError:
            raise ValueError(f"python-docx not installed, cannot extract {filename}")

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            slides_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_parts.append(para.text)
                if slide_parts:
                    slides_text.append(f"## Slide {slide_num}\n\n" + "\n".join(slide_parts))
            text = "\n\n".join(slides_text)
            if not text.strip():
                raise ValueError(f"No text extracted from {filename}")
            logger.info(f"python-pptx extracted {len(text)} chars from {filename}")
            return text
        except ImportError:
            raise ValueError(f"python-pptx not installed, cannot extract {filename}")

    raise ValueError(f"Unsupported office format: {ext}")


def _parse_url_list(urls_json: Optional[str]) -> List[str]:
    """Parse a JSON string into a list of URL strings."""
    if not urls_json:
        return []
    try:
        urls = json.loads(urls_json)
        if not isinstance(urls, list):
            raise ValueError("Expected a JSON array of URLs")
        return [str(u) for u in urls]
    except json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON for URLs: {e}")


async def _process_youtube_urls(urls: List[str], model: Optional[str] = None) -> List[dict]:
    """Process YouTube URLs by extracting transcripts (parallel)."""
    semaphore = asyncio.Semaphore(4)

    async def _process_single_youtube(url):
        async with semaphore:
            try:
                video_id = yt_service.extract_video_id(url)
                if not video_id:
                    logger.warning(f"Invalid YouTube URL, could not extract video ID: {url}")
                    return None

                result = await asyncio.to_thread(yt_service.get_transcript, url)
                if not result.get("success"):
                    logger.warning(f"YouTube transcript failed for {url}: {result.get('message', 'Unknown error')}")
                    return None

                text = result.get("text", "")
                if len(text) < 50:
                    logger.warning(f"YouTube transcript too short for {url}: {len(text)} chars")
                    return None

                if "video_title" in result and result["video_title"]:
                    video_title = result["video_title"]
                else:
                    video_title = await asyncio.to_thread(yt_service.get_video_title, url)

                if not video_title or video_title.startswith("YouTube Video:") or len(video_title) < 5:
                    video_title = await _extract_topic(text[:2000], model=model)

                logger.info(f"Processed YouTube {video_id}: {video_title} ({len(text)} chars, method={result.get('method', 'unknown')})")
                return {
                    "filename": f"youtube_{video_id}.txt",
                    "topic": video_title, "extracted_text": text,
                    "pages": 1, "char_count": len(text),
                    "source_url": url, "source_type": "youtube",
                }
            except Exception as e:
                logger.warning(f"YouTube processing error for {url}: {e}")
                return None

    results = await asyncio.gather(*[_process_single_youtube(url) for url in urls])
    return [r for r in results if r is not None]


async def _process_web_urls(urls: List[str], model: Optional[str] = None) -> List[dict]:
    """Process web URLs by extracting content via Jina Reader (parallel)."""
    semaphore = asyncio.Semaphore(4)

    async def _process_single_web_url(url):
        async with semaphore:
            try:
                result = await asyncio.to_thread(extract_web_content, url)
                if not result.get("success"):
                    logger.warning(f"Web extraction failed for {url}: {result.get('message', 'Unknown error')}")
                    return None

                text = result.get("text", "")
                if len(text) < 50:
                    logger.warning(f"Web content too short for {url}: {len(text)} chars")
                    return None

                topic = await _extract_topic(text[:2000], model=model)
                safe_domain = re.sub(r"[^a-zA-Z0-9]", "_", url.split("//")[-1].split("/")[0])

                logger.info(f"Processed web URL {url}: {topic} ({len(text)} chars)")
                return {
                    "filename": f"web_{safe_domain}.txt",
                    "topic": topic, "extracted_text": text,
                    "pages": 1, "char_count": len(text),
                    "source_url": url, "source_type": "web",
                }
            except Exception as e:
                logger.warning(f"Web URL processing error for {url}: {e}")
                return None

    results = await asyncio.gather(*[_process_single_web_url(url) for url in urls])
    return [r for r in results if r is not None]


def _save_temp_file(file: UploadFile, content: bytes = None) -> str:
    """Save uploaded file to temp location. Uses pre-read content if provided."""
    suffix = os.path.splitext(file.filename)[1] if file.filename else ""
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        data = content if content is not None else file.file.read()
        if not data:
            raise ValueError(f"Empty file: {file.filename}")
        tmp.write(data)
        return tmp.name


async def _extract_topic(text: str, model: Optional[str] = None) -> str:
    """Extract main topic from text using LLM"""
    prompt = build_topic_extraction_prompt(text)
    model_config = ModelConfig.get_config(model)

    # Use module-level singleton instead of creating new instances
    response = await course_service._call_model(prompt, model_config, expect_json=False)
    topic = response.get("content", "").strip()
    
    # Clean up
    topic = topic.replace('"', '').replace("Topic:", "").strip()
    
    return topic or "General Topic"


# Model info endpoint
@router.get("/models")
async def get_available_models():
    """List all available AI models for course generation"""
    from utils.model_config import MODEL_CONFIGS, DEFAULT_MODEL, estimate_course_cost, get_search_mode

    models = []
    for key, config in MODEL_CONFIGS.items():
        models.append({
            "id": key,
            "name": config["model"],
            "provider": config["provider"],
            "supports_web_search": config["supports_search"],
            "search_mode": get_search_mode(key),
            "estimated_cost_per_course": estimate_course_cost(key, 4),
            "default": key == DEFAULT_MODEL
        })

    return {"models": models}
