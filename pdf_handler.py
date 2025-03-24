import io
import re
from datetime import datetime
import PyPDF2
from fastapi import HTTPException
import httpx
from langchain.text_splitter import RecursiveCharacterTextSplitter
import chromadb
import logging
import os
from pptx import Presentation
from io import BytesIO
from docx import Document

logger = logging.getLogger(__name__)

class PDFHandler:
    def __init__(self):
        # Initialize ChromaDB client
        chroma_host = os.getenv('CHROMA_HOST', os.getenv('CHROMA_DB_CONNECTION_STRING'))
        chroma_port = int(os.getenv('CHROMA_PORT', 8000))
        self.chroma_client = chromadb.HttpClient(
            host=chroma_host, 
            port=chroma_port,
        )
    
    @staticmethod
    def clean_filename(filename):
        """
        Cleans the filename by replacing special characters and spaces with underscores.
        Converts the filename to lowercase.
        """
        cleaned_name = re.sub(r'[^a-zA-Z0-9._-]', '_', filename)
        cleaned_name = cleaned_name.lower()
        cleaned_name = re.sub(r'_+', '_', cleaned_name)
        return cleaned_name

    @staticmethod
    def determine_page_for_chunk(chunk: str, page_text_map: list[dict], key: str) -> int:
        """
        Determine the page or slide number for a given text chunk.
        
        Args:
        - chunk: Text chunk to locate
        - page_text_map: List of page or slide text mappings
        - key: The key to use for determining the page/slide number ('page_num' or 'slide_num')
        
        Returns:
        - Page or slide number where the chunk is most likely located
        """
        for page_info in page_text_map:
            if chunk in page_info['text']:
                return page_info[key]  # Use the specified key to get the number
        return page_text_map[-1][key]  # Return the last page/slide number if not found

    @staticmethod
    def extract_pptx_text(file_content):
        """
        Extracts text from a PowerPoint (.pptx) file.
        Returns a tuple of (full_text, slide_text_map) where slide_text_map contains text per slide.
        """
        full_text = ""
        slide_text_map = []
        
        # Load the presentation from bytes
        prs = Presentation(BytesIO(file_content))
        
        # Iterate through all slides
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = ""
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text += text + "\n"
            
            # Clean the extracted text
            slide_text = ''.join(char for char in slide_text if char.isprintable())
            
            # Add to the slide map
            slide_text_map.append({
                'start': len(full_text),
                'text': slide_text,
                'slide_num': slide_num
            })
            
            full_text += slide_text + "\n"
        
        return full_text.strip(), slide_text_map

    @staticmethod
    def extract_docx_text(file_content):
        """
        Extracts text from a Word (.docx) file.
        Returns a tuple of (full_text, paragraph_text_map).
        """
        full_text = ""
        paragraph_text_map = []
        
        # Load the document from bytes
        doc = Document(BytesIO(file_content))
        
        # Extract text from paragraphs
        for para_num, paragraph in enumerate(doc.paragraphs, start=1):
            text = paragraph.text.strip()
            if text:
                # Clean the extracted text
                text = ''.join(char for char in text if char.isprintable())
                
                # Add to the paragraph map
                paragraph_text_map.append({
                    'start': len(full_text),
                    'text': text,
                    'paragraph_num': para_num
                })
                
                full_text += text + "\n"
        
        return full_text.strip(), paragraph_text_map

    async def handle_pdf_upload(self, file, space_id, rag_system, nuton_api, file_path):
        try:
            # Read PDF
            pdf_content = await file.read()
            
            try:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_content))
            except Exception as pdf_error:
                logger.error(f"Error reading PDF: {pdf_error}")
                raise ValueError(f"Invalid or corrupted PDF file: {str(pdf_error)}")
            
            # Extract text with advanced page tracking
            full_text = ""
            page_text_map = []
            
            for page_num, page in enumerate(pdf_reader.pages, start=1):
                try:
                    try:
                        page_content = page.extract_text()
                    except Exception:
                        try:
                            page_content = page.get_text()
                        except Exception as alternative_extract_error:
                            logger.warning(f"Failed to extract text from page {page_num}: {alternative_extract_error}")
                            page_content = ""
                    
                    page_content = ''.join(char for char in page_content if char.isprintable())
                    
                    page_text_map.append({
                        'start': len(full_text),
                        'text': page_content,
                        'page_num': page_num
                    })
                    full_text += page_content + "\n"
                
                except Exception as page_error:
                    logger.error(f"Unexpected error extracting text from page {page_num}: {page_error}")
                    continue
            
            if not full_text.strip():
                raise ValueError("No text could be extracted from the PDF. The file may be scanned, encrypted, or contain only images.")
            
            print("space_id", space_id)
            # Save PDF metadata to database
            try:
                pdf_upload = rag_system.supabase.table('pdfs').insert({
                    'space_id': space_id,
                    'extracted_text': full_text,
                    'file_type': 'pdf',
                    'file_path': file_path
                }).execute()
                
                pdf_id = pdf_upload.data[0]['id']
            except Exception as db_error:
                logger.error(f"Database error: {db_error}")
                raise ValueError(f"Failed to save PDF metadata to database: {str(db_error)}")

            # Trigger upload in the background
            try:
                async with httpx.AsyncClient() as client:
                    response = await client.post(
                        f"{nuton_api}/upload",
                        files={"file": (file.filename, pdf_content, file.content_type)},
                        data={"pdf_id": pdf_id},
                        timeout=60.0
                    )

                    if response.status_code != 200:
                        logger.error(f"Error from worker service: {response.text}")
                        # Continue processing even if worker service fails
            except Exception as upload_error:
                logger.error(f"Error uploading to worker service: {upload_error}")
                # Continue processing even if worker service fails
            
            # Process text chunks and embeddings
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=10000,
                chunk_overlap=100
            )
            text_chunks = text_splitter.split_text(full_text)
            
            try:
                collection = self.chroma_client.get_or_create_collection("pdf_embeddings")
            except Exception as chroma_error:
                logger.error(f"ChromaDB error: {chroma_error}")
                raise ValueError(f"Failed to connect to vector database: {str(chroma_error)}")
            
            chunk_errors = 0
            for i, chunk in enumerate(text_chunks):
                try:
                    page_number = self.determine_page_for_chunk(chunk, page_text_map, 'page_num')
                    embedding = rag_system.generate_embedding(chunk)
                    
                    collection.add(
                        ids=[f"{pdf_id}_{i}"],
                        embeddings=[embedding],
                        metadatas=[{
                            "pdf_id": pdf_id,
                            "chunk_index": i,
                            "page": page_number,
                            "created_at": datetime.now().isoformat()
                        }],
                        documents=[chunk]
                    )
                except Exception as chunk_error:
                    logger.error(f"Error processing chunk {i}: {chunk_error}")
                    chunk_errors += 1
                    # Continue processing other chunks
            
            if chunk_errors == len(text_chunks):
                raise ValueError("Failed to process any text chunks. Check embedding service.")
            
            return {
                "status": "success", 
                "pdf_id": pdf_id,
                "extracted_text": full_text, 
                "chunks_processed": len(text_chunks) - chunk_errors,
                "chunks_failed": chunk_errors
            }
        
        except ValueError as ve:
            # Re-raise ValueError as HTTPException with 400 status code
            raise HTTPException(status_code=400, detail=str(ve))
        except Exception as e:
            logger.error(f"Unexpected error in handle_pdf_upload: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to process PDF: {str(e)}")

    async def handle_pptx_upload(self, file, space_id, rag_system, nuton_api, file_path):
        """
        Handles the upload and processing of PowerPoint (.pptx) files.
        Similar to handle_pdf_upload but for PowerPoint presentations.
        """
        try:
            # Read PPTX
            pptx_content = await file.read()
            
            try:
                full_text, slide_text_map = self.extract_pptx_text(pptx_content)
            except Exception as pptx_error:
                logger.error(f"Error extracting text from PPTX: {pptx_error}")
                raise ValueError(f"Invalid or corrupted PowerPoint file: {str(pptx_error)}")
            
            if not full_text.strip():
                raise ValueError("No text could be extracted from the PowerPoint file. The file may contain only images or non-textual content.")
            
            # Save PPTX metadata to database
            try:
                pptx_upload = rag_system.supabase.table('pdfs').insert({
                    'space_id': space_id,
                    'extracted_text': full_text,
                    'file_type': 'pptx',
                    'file_path': file_path
                }).execute()
                
                pptx_id = pptx_upload.data[0]['id']
            except Exception as db_error:
                logger.error(f"Database error: {db_error}")
                raise ValueError(f"Failed to save PowerPoint metadata to database: {str(db_error)}")

            # Trigger upload in the background
            try:
                async with httpx.AsyncClient() as client:
                    response = await client.post(
                        f"{nuton_api}/upload",
                        files={"file": (file.filename, pptx_content, file.content_type)},
                        data={"pdf_id": pptx_id},  # Using pdf_id for consistency with existing structure
                        timeout=60.0
                    )

                    if response.status_code != 200:
                        logger.error(f"Error from worker service: {response.text}")
                        # Continue processing even if worker service fails
            except Exception as upload_error:
                logger.error(f"Error uploading to worker service: {upload_error}")
                # Continue processing even if worker service fails
            
            # Process text chunks and embeddings
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=10000,
                chunk_overlap=100
            )
            text_chunks = text_splitter.split_text(full_text)
            
            try:
                collection = self.chroma_client.get_or_create_collection("pdf_embeddings")
            except Exception as chroma_error:
                logger.error(f"ChromaDB error: {chroma_error}")
                raise ValueError(f"Failed to connect to vector database: {str(chroma_error)}")

            chunk_errors = 0
            for i, chunk in enumerate(text_chunks):
                try:
                    # Find the slide number for this chunk using 'slide_num' key
                    slide_number = self.determine_page_for_chunk(chunk, slide_text_map, 'slide_num')
                    embedding = rag_system.generate_embedding(chunk)
                    
                    collection.add(
                        ids=[f"{pptx_id}_{i}"],
                        embeddings=[embedding],
                        metadatas=[{
                            "pdf_id": pptx_id,  # Using pdf_id for consistency
                            "chunk_index": i,
                            "page": slide_number,  # This will represent the slide number
                            "created_at": datetime.now().isoformat()
                        }],
                        documents=[chunk]
                    )
                except Exception as chunk_error:
                    logger.error(f"Error processing chunk {i}: {chunk_error}")
                    chunk_errors += 1
                    # Continue processing other chunks
            
            if chunk_errors == len(text_chunks):
                raise ValueError("Failed to process any text chunks. Check embedding service.")
            
            return {
                "status": "success", 
                "pdf_id": pptx_id,  # Using pdf_id for consistency
                "extracted_text": full_text, 
                "chunks_processed": len(text_chunks) - chunk_errors,
                "chunks_failed": chunk_errors
            }
        
        except ValueError as ve:
            # Re-raise ValueError as HTTPException with 400 status code
            raise HTTPException(status_code=400, detail=str(ve))
        except Exception as e:
            logger.error(f"Unexpected error in handle_pptx_upload: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to process PowerPoint file: {str(e)}")

    async def handle_docx_upload(self, file, space_id, rag_system, nuton_api, file_path):
        """
        Handles the upload and processing of Word (.docx) files.
        Similar to handle_pdf_upload but for Word documents.
        """
        try:
            # Read DOCX
            docx_content = await file.read()
            
            try:
                full_text, paragraph_text_map = self.extract_docx_text(docx_content)
            except Exception as docx_error:
                logger.error(f"Error extracting text from DOCX: {docx_error}")
                raise ValueError(f"Invalid or corrupted Word document: {str(docx_error)}")
            
            if not full_text.strip():
                raise ValueError("No text could be extracted from the Word document. The file may contain only images or non-textual content.")
            
            # Save DOCX metadata to database
            try:
                docx_upload = rag_system.supabase.table('pdfs').insert({
                    'space_id': space_id,
                    'extracted_text': full_text,
                    'file_type': 'docx',
                    'file_path': file_path
                }).execute()
                
                docx_id = docx_upload.data[0]['id']
            except Exception as db_error:
                logger.error(f"Database error: {db_error}")
                raise ValueError(f"Failed to save Word document metadata to database: {str(db_error)}")

            # Trigger upload in the background
            try:
                async with httpx.AsyncClient() as client:
                    response = await client.post(
                        f"{nuton_api}/upload",
                        files={"file": (file.filename, docx_content, file.content_type)},
                        data={"pdf_id": docx_id},  # Using pdf_id for consistency
                        timeout=60.0
                    )

                    if response.status_code != 200:
                        logger.error(f"Error from worker service: {response.text}")
                        # Continue processing even if worker service fails
            except Exception as upload_error:
                logger.error(f"Error uploading to worker service: {upload_error}")
                # Continue processing even if worker service fails
            
            # Process text chunks and embeddings
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=10000,
                chunk_overlap=100
            )
            text_chunks = text_splitter.split_text(full_text)
            
            try:
                collection = self.chroma_client.get_or_create_collection("pdf_embeddings")
            except Exception as chroma_error:
                logger.error(f"ChromaDB error: {chroma_error}")
                raise ValueError(f"Failed to connect to vector database: {str(chroma_error)}")
            
            chunk_errors = 0
            for i, chunk in enumerate(text_chunks):
                try:
                    # Find the paragraph number for this chunk
                    paragraph_number = self.determine_page_for_chunk(chunk, paragraph_text_map, 'paragraph_num')
                    embedding = rag_system.generate_embedding(chunk)
                    
                    collection.add(
                        ids=[f"{docx_id}_{i}"],
                        embeddings=[embedding],
                        metadatas=[{
                            "pdf_id": docx_id,  # Using pdf_id for consistency
                            "chunk_index": i,
                            "page": paragraph_number,  # This will represent the paragraph number
                            "created_at": datetime.now().isoformat()
                        }],
                        documents=[chunk]
                    )
                except Exception as chunk_error:
                    logger.error(f"Error processing chunk {i}: {chunk_error}")
                    chunk_errors += 1
                    # Continue processing other chunks
            
            if chunk_errors == len(text_chunks):
                raise ValueError("Failed to process any text chunks. Check embedding service.")
            
            return {
                "status": "success", 
                "pdf_id": docx_id,  # Using pdf_id for consistency
                "extracted_text": full_text, 
                "chunks_processed": len(text_chunks) - chunk_errors,
                "chunks_failed": chunk_errors
            }
        
        except ValueError as ve:
            # Re-raise ValueError as HTTPException with 400 status code
            raise HTTPException(status_code=400, detail=str(ve))
        except Exception as e:
            logger.error(f"Unexpected error in handle_docx_upload: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to process Word document: {str(e)}") 